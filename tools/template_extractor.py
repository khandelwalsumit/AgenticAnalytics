"""Extract slide layout structure from a PPTX template into a JSON catalog.

Run once offline whenever the template changes:
    python tools/template_extractor.py

Produces ``data/input/template_catalog.json`` with:
- Every slide layout: name, index, placeholder positions/sizes/types
- A curated ``section_map`` that assigns specific layouts to each report section
- A ``visual_hierarchy`` block defining font sizes per heading level

The formatting agent and the deterministic PPTX builder both consume this catalog
so they share the same source-of-truth for what each slide can contain.
"""

from __future__ import annotations

import json
import sys
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Emu


# ---------------------------------------------------------------------------
# Placeholder type enum → human-readable string
# ---------------------------------------------------------------------------

_PH_TYPE_NAMES = {
    1: "TITLE",
    2: "BODY",
    3: "CENTER_TITLE",
    4: "SUBTITLE",
    7: "OBJECT",
    13: "SLIDE_NUMBER",
    15: "FOOTER",
    18: "PICTURE",
}


def _ph_type_name(ph_type_value: int) -> str:
    return _PH_TYPE_NAMES.get(ph_type_value, f"UNKNOWN({ph_type_value})")


def _emu_to_inches(emu: int) -> float:
    return round(emu / 914400, 2)


# ---------------------------------------------------------------------------
# Core extractor
# ---------------------------------------------------------------------------


def extract_layouts(template_path: str) -> dict[str, Any]:
    """Read a .pptx and return the full layout catalog as a dict."""
    prs = Presentation(template_path)

    catalog: dict[str, Any] = {
        "template_file": Path(template_path).name,
        "slide_width_inches": _emu_to_inches(prs.slide_width),
        "slide_height_inches": _emu_to_inches(prs.slide_height),
        "layouts": {},
    }

    for idx, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for ph in layout.placeholders:
            ph_type = ph.placeholder_format.type
            # skip SLIDE_NUMBER and FOOTER — not content-relevant
            if ph_type in (13, 15):
                continue
            placeholders.append({
                "idx": ph.placeholder_format.idx,
                "name": ph.name,
                "type": _ph_type_name(ph_type),
                "left_in": _emu_to_inches(ph.left),
                "top_in": _emu_to_inches(ph.top),
                "width_in": _emu_to_inches(ph.width),
                "height_in": _emu_to_inches(ph.height),
            })

        catalog["layouts"][str(idx)] = {
            "index": idx,
            "name": layout.name,
            "placeholders": placeholders,
        }

    # Add curated section map and visual hierarchy
    catalog["section_map"] = _build_section_map(catalog["layouts"])
    catalog["visual_hierarchy"] = _build_visual_hierarchy()

    return catalog


# ---------------------------------------------------------------------------
# Section map — which template layout to use for each slide type
# ---------------------------------------------------------------------------


def _build_section_map(layouts: dict[str, Any]) -> dict[str, Any]:
    """Map report sections to the best template layout by name matching."""

    def _find_layout_index(name_prefix: str) -> int | None:
        for key, layout in layouts.items():
            if layout["name"] == name_prefix:
                return layout["index"]
        return None

    # Fallbacks: try exact name, then first match
    title_idx = _find_layout_index("Title") or 6
    content_idx = _find_layout_index("Content") or 1
    section_header_idx = _find_layout_index("Section Header") or 9
    number_large_idx = _find_layout_index("Number Large") or 37
    title_only_idx = _find_layout_index("Title Only") or 51
    content_with_pic_idx = _find_layout_index("Content with Picture 2") or 19
    two_content_idx = _find_layout_index("Two Content") or 47
    conclusion_idx = _find_layout_index("Conclusion") or 55

    return {
        "exec_summary": {
            "description": "Executive Summary \u2014 2 slides",
            "slides": [
                {
                    "slide_role": "executive_summary",
                    "layout_index": title_idx,
                    "layout_name": layouts.get(str(title_idx), {}).get("name", "Title"),
                    "placeholders_used": ["TITLE", "SUBTITLE"],
                    "content_guidance": "Title: 'EXECUTIVE SUMMARY'. Subtitle: context sentence. Quick Wins: 3 action items with call impact.",
                    "structured_fields": ["title", "subtitle", "quick_wins"],
                },
                {
                    "slide_role": "pain_points",
                    "layout_index": content_idx,
                    "layout_name": layouts.get(str(content_idx), {}).get("name", "Content"),
                    "placeholders_used": ["TITLE", "OBJECT"],
                    "content_guidance": "Title: assertion with numbers. 3 structured pain point cards, each with 2-3 line issue and 1-2 line fix with owner in parens.",
                    "card_count": 3,
                    "structured_fields": ["title", "cards"],
                    "card_fields": ["name", "calls", "impact_score", "priority", "issue", "fix"],
                },
            ],
        },
        "impact": {
            "description": "Impact & Prioritization \u2014 3 slides",
            "slides": [
                {
                    "slide_role": "impact_matrix",
                    "layout_index": title_only_idx,
                    "layout_name": layouts.get(str(title_only_idx), {}).get("name", "Title Only"),
                    "placeholders_used": ["TITLE"],
                    "content_guidance": "Theme card list LEFT (~60%), scatter chart RIGHT (~40%). Each theme: name + quadrant, stats, issue. NOT a table.",
                    "structured_fields": ["title", "themes", "chart_placeholder"],
                },
                {
                    "slide_role": "low_hanging_fruit",
                    "layout_index": content_idx,
                    "layout_name": layouts.get(str(content_idx), {}).get("name", "Content"),
                    "placeholders_used": ["TITLE", "OBJECT"],
                    "content_guidance": "3 easiest-to-implement solutions sorted by ease. Each: title (blue 16pt), detail (12pt elaboration), call_impact.",
                    "structured_fields": ["title", "solutions"],
                    "solution_fields": ["title", "detail", "call_impact"],
                },
                {
                    "slide_role": "recommendations",
                    "layout_index": content_idx,
                    "layout_name": layouts.get(str(content_idx), {}).get("name", "Content"),
                    "placeholders_used": ["TITLE", "OBJECT"],
                    "content_guidance": "2x2 grid of dimension cards. Each dimension has name, accent_color, and 1-2 consolidated actions.",
                    "structured_fields": ["title", "dimensions"],
                    "dimension_fields": ["name", "accent_color", "actions"],
                },
            ],

        },
        "theme_deep_dives": {
            "description": "Theme Deep Dives — 1 slide per theme (max 10)",
            "per_theme_slide": {
                "slide_role": "theme_card",
                "layout_index": content_with_pic_idx,
                "layout_name": layouts.get(str(content_with_pic_idx), {}).get("name", "Content with Picture 2"),
                "placeholders_used": ["TITLE", "OBJECT", "PICTURE"],
                "content_guidance": "Two-column layout. LEFT (60%): core_issue + primary_driver + solutions. RIGHT (40%): driver_table. Stats bar below title shows calls/pct/impact/ease/priority.",
                "structured_fields": ["title", "stats_bar", "left_column", "right_column"],
                "left_column_fields": ["core_issue", "primary_driver", "solutions"],
                "right_column_fields": ["type", "headers", "rows"],
            },
            "max_themes": 10,
        },
    }


# ---------------------------------------------------------------------------
# Visual hierarchy — font specs for deterministic rendering
# ---------------------------------------------------------------------------


def _build_visual_hierarchy() -> dict[str, Any]:
    """McKinsey-style visual hierarchy for the PPTX builder."""
    return {
        "h1": {
            "font_name": "Calibri",
            "font_size_pt": 28,
            "bold": True,
            "color_hex": "003B70",
            "usage": "Slide title — assertion with data",
        },
        "h2": {
            "font_name": "Calibri",
            "font_size_pt": 20,
            "bold": True,
            "color_hex": "003B70",
            "usage": "Section sub-heading within slide body",
        },
        "h3": {
            "font_name": "Calibri",
            "font_size_pt": 16,
            "bold": True,
            "color_hex": "333333",
            "usage": "Dimension label or group heading (e.g., Digital / UX)",
        },
        "point_heading": {
            "font_name": "Calibri",
            "font_size_pt": 14,
            "bold": True,
            "color_hex": "333333",
            "usage": "Bold label before a bullet (e.g., 'What\\'s happening:')",
        },
        "point_description": {
            "font_name": "Calibri",
            "font_size_pt": 13,
            "bold": False,
            "color_hex": "333333",
            "usage": "Normal bullet body text",
        },
        "sub_point": {
            "font_name": "Calibri",
            "font_size_pt": 12,
            "bold": False,
            "color_hex": "666666",
            "usage": "Indented sub-bullet or secondary detail",
        },
        "callout": {
            "font_name": "Calibri",
            "font_size_pt": 24,
            "bold": True,
            "color_hex": "006BA6",
            "usage": "Big stat or key assertion (e.g., biggest bet callout)",
        },
        "table_header": {
            "font_name": "Calibri",
            "font_size_pt": 11,
            "bold": True,
            "color_hex": "FFFFFF",
            "bg_color_hex": "003B70",
            "usage": "Table column headers",
        },
        "table_cell": {
            "font_name": "Calibri",
            "font_size_pt": 10,
            "bold": False,
            "color_hex": "333333",
            "usage": "Table body cells",
        },
    }


# ---------------------------------------------------------------------------
# CLI entry point
# ---------------------------------------------------------------------------


def main() -> None:
    root = Path(__file__).resolve().parent.parent
    template_path = root / "data" / "input" / "template.pptx"
    output_path = root / "data" / "input" / "template_catalog.json"

    if not template_path.exists():
        print(f"ERROR: Template not found at {template_path}")
        sys.exit(1)

    catalog = extract_layouts(str(template_path))

    # Write pretty-printed JSON
    output_path.write_text(json.dumps(catalog, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"Template catalog written to {output_path}")
    print(f"  Layouts extracted: {len(catalog['layouts'])}")
    print(f"  Sections mapped: {list(catalog['section_map'].keys())}")


if __name__ == "__main__":
    main()
