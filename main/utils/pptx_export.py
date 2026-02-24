"""PPTX generator — template-based with code fallback.

Supports two modes:
1. External template: loads a .pptx template and uses its slide layouts
2. Code-based: creates a clean presentation with Citi-style defaults

The slide plan (from the Narrative Agent) drives which slides are created
and what content goes on each.
"""

from __future__ import annotations

import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor


# -- Layout indices for code-based template -----------------------------------
# These map to the default slide layouts in python-pptx's blank presentation.
# When using an external template, layout names are matched by convention.
LAYOUT_TITLE = 0        # Title Slide
LAYOUT_CONTENT = 1      # Title and Content
LAYOUT_BLANK = 6         # Blank


# -- Brand defaults (used when no external template is provided) --------------
BRAND = {
    "title_font": "Calibri",
    "body_font": "Calibri",
    "title_size": Pt(28),
    "subtitle_size": Pt(18),
    "body_size": Pt(14),
    "bullet_size": Pt(13),
    "title_color": RGBColor(0x00, 0x3B, 0x70),    # Dark blue
    "subtitle_color": RGBColor(0x66, 0x66, 0x66),  # Grey
    "body_color": RGBColor(0x33, 0x33, 0x33),       # Dark grey
    "accent_color": RGBColor(0x00, 0x6B, 0xA6),     # Citi blue
}


def _find_layout(prs: Presentation, layout_type: str) -> Any:
    """Find a slide layout by name convention or fall back to index."""
    name_map = {
        "title": ["Title Slide", "Title", "title"],
        "content": ["Title and Content", "Content", "Title, Content"],
        "blank": ["Blank", "blank"],
    }
    target_names = name_map.get(layout_type, [])

    for layout in prs.slide_layouts:
        if layout.name in target_names:
            return layout

    # Fallback to index
    idx_map = {"title": LAYOUT_TITLE, "content": LAYOUT_CONTENT, "blank": LAYOUT_BLANK}
    idx = idx_map.get(layout_type, LAYOUT_CONTENT)
    if idx < len(prs.slide_layouts):
        return prs.slide_layouts[idx]
    return prs.slide_layouts[0]


def _strip_markdown(text: str) -> str:
    """Remove markdown bold/italic markers for clean PPTX text."""
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*", r"\1", text)
    text = re.sub(r"`(.*?)`", r"\1", text)
    return text


def _add_title_slide(prs: Presentation, slide_data: dict) -> None:
    """Add a title slide with title and subtitle."""
    layout = _find_layout(prs, "title")
    slide = prs.slides.add_slide(layout)

    if slide.shapes.title:
        slide.shapes.title.text = _strip_markdown(slide_data.get("title", ""))
        for p in slide.shapes.title.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = BRAND["title_size"]
                run.font.color.rgb = BRAND["title_color"]
                run.font.bold = True

    # Find subtitle placeholder
    subtitle = slide_data.get("subtitle", "")
    if subtitle and len(slide.placeholders) > 1:
        ph = slide.placeholders[1]
        ph.text = _strip_markdown(subtitle)
        for p in ph.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = BRAND["subtitle_size"]
                run.font.color.rgb = BRAND["subtitle_color"]


def _add_content_slide(prs: Presentation, slide_data: dict) -> None:
    """Add a content slide with title and bullet points."""
    layout = _find_layout(prs, "content")
    slide = prs.slides.add_slide(layout)

    # Title
    if slide.shapes.title:
        slide.shapes.title.text = _strip_markdown(slide_data.get("title", ""))
        for p in slide.shapes.title.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = BRAND["title_size"]
                run.font.color.rgb = BRAND["title_color"]
                run.font.bold = True

    # Body — bullet points
    points = slide_data.get("points", [])
    if points and len(slide.placeholders) > 1:
        body_ph = slide.placeholders[1]
        tf = body_ph.text_frame
        tf.clear()

        for i, point in enumerate(points):
            text = _strip_markdown(point)

            # Determine indent level from prefix
            indent_level = 0
            if text.startswith("  - ") or text.startswith("  * "):
                text = text[4:]
                indent_level = 1
            elif text.startswith("- ") or text.startswith("* "):
                text = text[2:]

            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.text = text
            p.level = indent_level
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = BRAND["body_size"]
                run.font.color.rgb = BRAND["body_color"]


def _add_chart_slide(prs: Presentation, slide_data: dict, chart_paths: dict[str, str]) -> None:
    """Add a content slide with chart image embedded.

    If the chart image exists, places it on the right half of the slide.
    Points go on the left half.
    """
    layout = _find_layout(prs, "content")
    slide = prs.slides.add_slide(layout)

    # Title
    if slide.shapes.title:
        slide.shapes.title.text = _strip_markdown(slide_data.get("title", ""))
        for p in slide.shapes.title.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = BRAND["title_size"]
                run.font.color.rgb = BRAND["title_color"]
                run.font.bold = True

    # Bullet points in body placeholder
    points = slide_data.get("points", [])
    if points and len(slide.placeholders) > 1:
        body_ph = slide.placeholders[1]
        tf = body_ph.text_frame
        tf.clear()

        for i, point in enumerate(points):
            text = _strip_markdown(point)
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = text
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = BRAND["bullet_size"]
                run.font.color.rgb = BRAND["body_color"]

    # Embed chart image if available
    visual_id = slide_data.get("visual", "")
    if visual_id and visual_id != "none":
        chart_path = chart_paths.get(visual_id, "")
        if chart_path and Path(chart_path).exists():
            # Place chart on the right side of the slide
            left = Inches(6.5)
            top = Inches(1.8)
            width = Inches(6.0)
            slide.shapes.add_picture(chart_path, left, top, width=width)


def generate_pptx_from_slides(
    slide_plan: dict,
    chart_paths: dict[str, str],
    output_path: str,
    template_path: str = "",
) -> str:
    """Generate a PPTX from a structured slide plan.

    Args:
        slide_plan: Dict with "slides" list from the Narrative Agent.
                    Each slide has: type, title, points, visual, notes.
        chart_paths: Map of visual_id → chart image file path.
        output_path: Where to save the .pptx file.
        template_path: Optional path to .pptx template file.
                       If empty or not found, uses code-based defaults.

    Returns:
        The output path.
    """
    # Load template or create blank presentation
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # Widescreen 16:9
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    slides = slide_plan.get("slides", [])

    for slide_data in slides:
        slide_type = slide_data.get("type", "content")

        if slide_type == "title":
            _add_title_slide(prs, slide_data)

        elif slide_type in ("theme_detail",) and slide_data.get("visual", "none") != "none":
            _add_chart_slide(prs, slide_data, chart_paths)

        elif slide_type == "impact_ease" and slide_data.get("visual", "none") != "none":
            _add_chart_slide(prs, slide_data, chart_paths)

        else:
            # key_summary, recommendations, appendix, or any other type
            _add_content_slide(prs, slide_data)

    prs.save(output_path)
    return output_path


# Legacy compatibility — used by generate_markdown_report flow
def markdown_to_pptx(markdown: str, output_path: str) -> str:
    """Convert a markdown report to a PowerPoint presentation.

    Parses markdown headers (# / ##) as slide titles, bullet points as
    content. This is the fallback when no structured slide plan is available.

    Args:
        markdown: The markdown content string.
        output_path: File path to save the .pptx.

    Returns:
        The output path.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    sections = _split_sections(markdown)

    for section in sections:
        _add_content_slide(prs, {
            "title": section["title"],
            "points": [line for line in section["body"].split("\n") if line.strip() and line.strip() != "---"],
        })

    prs.save(output_path)
    return output_path


def _split_sections(markdown: str) -> list[dict]:
    """Split markdown into sections by ## headers."""
    sections = []
    current_title = ""
    current_body: list[str] = []

    for line in markdown.split("\n"):
        if line.startswith("# ") and not line.startswith("## "):
            if current_title or current_body:
                sections.append({"title": current_title, "body": "\n".join(current_body)})
            current_title = line.lstrip("# ").strip()
            current_body = []
        elif line.startswith("## "):
            if current_title or current_body:
                sections.append({"title": current_title, "body": "\n".join(current_body)})
            current_title = line.lstrip("# ").strip()
            current_body = []
        else:
            current_body.append(line)

    if current_title or current_body:
        sections.append({"title": current_title, "body": "\n".join(current_body)})

    return sections
