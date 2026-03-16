"""Deterministic PPTX builder — section-aware, template-based.

Takes per-section JSON blueprints (from the formatting agent) + template.pptx +
chart images and produces the final report.pptx.

No LLM involved — placement, fonts, and structure are fully deterministic.

Slide types supported:
  1. executive_summary — title + subtitle_lines + quick wins
  2. pain_points — 3-column card layout
  3. impact_ease — theme card list LEFT + scatter chart RIGHT
  4. low_hanging_fruit — 3 easiest solutions
  5. recommendations — 2x2 dimension grid
  6. theme_card — stats bar + two-column narrative/table

Canvas: 13.33 x 7.5 inches (widescreen 16:9)
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

logger = logging.getLogger("agenticanalytics.pptx_builder")

# ═══════════════════════════════════════════════════════════════════════════
# Global constants — from slide_improvement_plan.md
# ═══════════════════════════════════════════════════════════════════════════

SLIDE_W = 13.33
SLIDE_H = 7.5
MARGIN_L = 0.5
MARGIN_R = 0.5
CONTENT_W = 12.33  # 13.33 - 0.5 - 0.5

# -- Color palette --
NAVY        = RGBColor(0x00, 0x3B, 0x70)
BLUE_ACCENT = RGBColor(0x00, 0x6B, 0xA6)
BLACK       = RGBColor(0x22, 0x22, 0x22)
DARK_GRAY   = RGBColor(0x44, 0x44, 0x44)
MID_GRAY    = RGBColor(0x88, 0x88, 0x88)
LIGHT_GRAY  = RGBColor(0xD0, 0xD0, 0xD0)
BG_LIGHT    = RGBColor(0xF5, 0xF7, 0xFA)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

ACCENT_RED    = RGBColor(0xC0, 0x39, 0x2B)
ACCENT_AMBER  = RGBColor(0xE6, 0x7E, 0x22)
ACCENT_GREEN  = RGBColor(0x27, 0xAE, 0x60)
ACCENT_PURPLE = RGBColor(0x8E, 0x44, 0xAD)

DIMENSION_COLORS = {
    "Digital / UX":   BLUE_ACCENT,
    "Digital/UX":     BLUE_ACCENT,
    "Operations":     ACCENT_GREEN,
    "Communications": ACCENT_AMBER,
    "Policy":         ACCENT_PURPLE,
}

FONT = "Calibri"


# -- Typography presets --
def _style(size, bold=False, color=BLACK):
    return {"size": Pt(size), "bold": bold, "color": color, "font": FONT}

SLIDE_TITLE  = _style(22, bold=True,  color=NAVY)
SUBTITLE_TEXT = _style(13, bold=False, color=DARK_GRAY)
H3_LABEL     = _style(14, bold=True,  color=BLUE_ACCENT)
BODY_TEXT     = _style(11, bold=False, color=BLACK)
BOLD_BODY    = _style(11, bold=True,  color=BLACK)
STATS_TEXT   = _style(10, bold=False, color=MID_GRAY)
SMALL_MUTED  = _style(9,  bold=False, color=MID_GRAY)
TABLE_HEADER = _style(10, bold=True,  color=WHITE)
TABLE_CELL   = _style(9,  bold=False, color=BLACK)
BULLET_TITLE = _style(12, bold=True,  color=BLUE_ACCENT)
BULLET_BODY  = _style(11, bold=False, color=BLACK)


# ═══════════════════════════════════════════════════════════════════════════
# Utility helpers
# ═══════════════════════════════════════════════════════════════════════════

def _strip_md(text: str) -> str:
    text = re.sub(r"\*\*\*(.*?)\*\*\*", r"\1", text)
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)
    text = re.sub(r"\*(.*?)\*", r"\1", text)
    text = re.sub(r"`(.*?)`", r"\1", text)
    return text.strip()


def _apply(run, style: dict):
    run.font.size = style["size"]
    run.font.bold = style["bold"]
    run.font.color.rgb = style["color"]
    run.font.name = style["font"]


def _hex_to_rgb(hex_str: str) -> RGBColor:
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


def _tb(slide, x, y, w, h, text, style, align=PP_ALIGN.LEFT, wrap=True):
    """Add a textbox with a single styled run."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = _strip_md(str(text))
    _apply(run, style)
    return box


def _rule(slide, x, y, w, color=LIGHT_GRAY):
    """Add a thin horizontal rectangle rule."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(0.01))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def _rect(slide, x, y, w, h, color):
    """Add a filled rectangle with no border."""
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _title_block(slide, title_text):
    """Shared title block: title at y=0.3, rule at y=0.72. Content starts at y=0.85."""
    _tb(slide, MARGIN_L, 0.3, CONTENT_W, 0.4, title_text, SLIDE_TITLE)
    _rule(slide, MARGIN_L, 0.72, CONTENT_W)


# ═══════════════════════════════════════════════════════════════════════════
# Table helper
# ═══════════════════════════════════════════════════════════════════════════

def _add_table(slide, headers, rows, x=0.5, y=1.0, w=9.0, col_widths=None, row_h=0.28):
    if not headers:
        return
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl = slide.shapes.add_table(n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(row_h * n_rows)).table

    # Column widths
    if col_widths and len(col_widths) == n_cols:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Inches(cw)

    # Header
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = _strip_md(str(h))
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            for r in p.runs:
                _apply(r, TABLE_HEADER)

    # Body
    for ri, row in enumerate(rows):
        for ci in range(n_cols):
            val = _strip_md(str(row[ci])) if ci < len(row) else ""
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_LIGHT
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                for r in p.runs:
                    _apply(r, TABLE_CELL)


def _add_chart_image(slide, chart_key, chart_paths, x=8.0, y=0.85, w=4.83):
    clean_key = chart_key
    m = re.search(r"\{\{\s*chart\.([a-zA-Z0-9_-]+)\s*\}\}", chart_key)
    if m:
        clean_key = m.group(1)
    path = chart_paths.get(clean_key, "")
    if path and Path(path).exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))


# ═══════════════════════════════════════════════════════════════════════════
# Slide renderers
# ═══════════════════════════════════════════════════════════════════════════


def _render_executive_summary(prs, slide_data, layout_idx):
    """Slide 1: EXECUTIVE SUMMARY — title + subtitle_lines + quick wins."""
    slide = prs.slides.add_slide(prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)])

    # 1. Title
    _tb(slide, MARGIN_L, 0.3, CONTENT_W, 0.4, slide_data.get("title", "EXECUTIVE SUMMARY"), SLIDE_TITLE)

    # 2. Rule below title
    _rule(slide, MARGIN_L, 0.72, CONTENT_W)

    # 3. Subtitle block — supports subtitle_lines (rich) or subtitle (plain string)
    subtitle_lines = slide_data.get("subtitle_lines", [])
    subtitle_plain = slide_data.get("subtitle", "")

    y = 0.85
    if subtitle_lines:
        box = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(y), Inches(CONTENT_W), Inches(0.85))
        tf = box.text_frame
        tf.word_wrap = True
        for li, line_obj in enumerate(subtitle_lines):
            line_text = _strip_md(str(line_obj.get("text", "") if isinstance(line_obj, dict) else line_obj))
            bold_part = line_obj.get("bold_part") if isinstance(line_obj, dict) else None
            p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
            p.space_after = Pt(2)
            if bold_part and bold_part in line_text:
                before, after = line_text.split(bold_part, 1)
                if before:
                    r = p.add_run()
                    r.text = before
                    _apply(r, SUBTITLE_TEXT)
                r_bold = p.add_run()
                r_bold.text = bold_part
                _apply(r_bold, _style(13, bold=True, color=BLUE_ACCENT))
                if after:
                    r_after = p.add_run()
                    r_after.text = after
                    _apply(r_after, SUBTITLE_TEXT)
            else:
                r = p.add_run()
                r.text = line_text
                _apply(r, SUBTITLE_TEXT)
    elif subtitle_plain:
        _tb(slide, MARGIN_L, y, CONTENT_W, 0.85, subtitle_plain, SUBTITLE_TEXT)

    # 4. Rule below subtitle
    _rule(slide, MARGIN_L, 1.72, CONTENT_W)

    # 5. "Quick Wins:" label
    quick_wins = slide_data.get("quick_wins", [])
    if quick_wins:
        _tb(slide, MARGIN_L, 1.88, CONTENT_W, 0.3, "Quick Wins:", H3_LABEL)

        # 6. Quick win items as rich text
        box = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(2.20), Inches(CONTENT_W), Inches(4.8))
        tf = box.text_frame
        tf.word_wrap = True
        for qi, qw in enumerate(quick_wins[:5]):
            if isinstance(qw, dict):
                action = _strip_md(str(qw.get("action", "")))
                detail = _strip_md(str(qw.get("detail", "")))
            else:
                # Plain string: split at " — " or use whole string
                action = _strip_md(str(qw))
                detail = ""

            p = tf.paragraphs[0] if qi == 0 else tf.add_paragraph()
            p.space_after = Pt(10)
            # Bold action title
            r_title = p.add_run()
            r_title.text = action
            _apply(r_title, BULLET_TITLE)
            # Detail on next line
            if detail:
                p_detail = tf.add_paragraph()
                r_det = p_detail.add_run()
                r_det.text = f"    {detail}"
                _apply(r_det, BODY_TEXT)


def _render_pain_points(prs, slide_data, layout_idx):
    """Slide 2: 3 pain point cards with accent bars."""
    slide = prs.slides.add_slide(prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)])

    _title_block(slide, slide_data.get("title", "Key Pain Points"))

    cards = slide_data.get("cards", [])
    accent_colors = [ACCENT_RED, ACCENT_AMBER, ACCENT_GREEN]

    # Card layout: 3 cards filling 13.33" canvas with 0.5" margins
    CARD_W = 3.77
    CARD_H = 5.5
    CARD_GAP = 0.35
    CARD_Y = 0.95
    CARD_X_START = (SLIDE_W - 3 * CARD_W - 2 * CARD_GAP) / 2  # ~0.66
    ACCENT_BAR_W = 0.06

    for i, card in enumerate(cards[:3]):
        card_x = CARD_X_START + i * (CARD_W + CARD_GAP)
        accent = accent_colors[i % 3]
        inner_x = card_x + 0.22
        inner_w = CARD_W - 0.37  # 3.40

        # Accent bar
        _rect(slide, card_x, CARD_Y, ACCENT_BAR_W, CARD_H, accent)
        # Card background
        _rect(slide, card_x + ACCENT_BAR_W, CARD_Y, CARD_W - ACCENT_BAR_W, CARD_H, BG_LIGHT)

        # Card title: "1. Theme Name (X calls)"
        name = _strip_md(str(card.get("name", f"Pain Point {i+1}")))
        calls = card.get("calls", 0)
        box = slide.shapes.add_textbox(Inches(inner_x), Inches(1.05), Inches(inner_w), Inches(0.40))
        tf = box.text_frame
        tf.word_wrap = True
        r1 = tf.paragraphs[0].add_run()
        r1.text = name
        _apply(r1, _style(11, bold=True, color=NAVY))
        r2 = tf.paragraphs[0].add_run()
        r2.text = f" ({calls} calls)" if calls else ""
        _apply(r2, _style(10, bold=False, color=NAVY))

        # Stats line
        impact = card.get("impact", card.get("impact_score", 0))
        priority = card.get("priority", 0)
        _tb(slide, inner_x, 1.45, inner_w, 0.25, f"Impact: {impact} | Priority: {priority}", _style(9, color=MID_GRAY))

        # Separator
        _rule(slide, inner_x, 1.72, inner_w)

        # Issue: label + text
        box = slide.shapes.add_textbox(Inches(inner_x), Inches(1.85), Inches(inner_w), Inches(1.6))
        tf = box.text_frame
        tf.word_wrap = True
        r_lbl = tf.paragraphs[0].add_run()
        r_lbl.text = "Issue: "
        _apply(r_lbl, _style(11, bold=True, color=BLACK))
        r_txt = tf.paragraphs[0].add_run()
        r_txt.text = _strip_md(str(card.get("issue", "")))
        _apply(r_txt, _style(10, color=BLACK))

        # Fix: label + text + (owner)
        box = slide.shapes.add_textbox(Inches(inner_x), Inches(3.55), Inches(inner_w), Inches(1.6))
        tf = box.text_frame
        tf.word_wrap = True
        r_fix_lbl = tf.paragraphs[0].add_run()
        r_fix_lbl.text = "Fix: "
        _apply(r_fix_lbl, _style(11, bold=True, color=BLACK))
        fix_text = _strip_md(str(card.get("fix", "")))
        r_fix = tf.paragraphs[0].add_run()
        r_fix.text = fix_text + " "
        _apply(r_fix, _style(10, color=BLACK))
        owner = card.get("owner", "")
        if owner and f"({owner})" not in fix_text:
            r_own = tf.paragraphs[0].add_run()
            r_own.text = f"({owner})"
            _apply(r_own, _style(9, bold=True, color=BLUE_ACCENT))


def _render_impact_ease(prs, slide_data, layout_idx, chart_paths):
    """Slide 3: Theme card list LEFT (58%) + scatter chart RIGHT (38%)."""
    slide = prs.slides.add_slide(prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)])

    _title_block(slide, slide_data.get("title", "Impact vs. Ease \u2014 Full Theme Prioritization"))

    themes = slide_data.get("themes", [])

    # Layout: LEFT 58% | gap | RIGHT 38%
    LEFT_W = 7.15
    RIGHT_X = MARGIN_L + LEFT_W + 0.3
    RIGHT_W = SLIDE_W - RIGHT_X - MARGIN_R  # ~4.88

    # LEFT SIDE — theme summary list
    if themes:
        box = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(0.95), Inches(LEFT_W), Inches(6.0))
        tf = box.text_frame
        tf.word_wrap = True

        for ti, t in enumerate(themes[:10]):
            t_name = _strip_md(str(t.get("name", "")))
            quadrant = _strip_md(str(t.get("quadrant", "")))
            # Derive quadrant if not provided
            if not quadrant:
                imp = float(t.get("impact", 0))
                ease = float(t.get("ease", 0))
                if imp >= 7 and ease >= 7:
                    quadrant = "Quick Win"
                elif imp >= 7:
                    quadrant = "Strategic Bet"
                elif ease >= 7:
                    quadrant = "Easy Fix"
                else:
                    quadrant = "Deprioritize"

            # Line 1: Bold theme name + quadrant
            p = tf.paragraphs[0] if ti == 0 else tf.add_paragraph()
            p.space_after = Pt(1)
            r_name = p.add_run()
            r_name.text = t_name
            _apply(r_name, _style(11, bold=True, color=NAVY))
            r_sep = p.add_run()
            r_sep.text = f" \u2014 {quadrant}"
            _apply(r_sep, _style(11, bold=False, color=MID_GRAY))

            # Line 2: Stats
            p2 = tf.add_paragraph()
            p2.space_after = Pt(6)
            r_stats = p2.add_run()
            r_stats.text = f"Calls: {t.get('calls', 0)} | Impact: {t.get('impact', 0)} | Ease: {t.get('ease', 0)} | Priority: {t.get('priority', 0)}"
            _apply(r_stats, _style(9, color=MID_GRAY))
    else:
        # Fallback: legacy table if present
        table_data = slide_data.get("table", {})
        if table_data and table_data.get("headers"):
            _add_table(slide, table_data["headers"], table_data.get("rows", []),
                       x=MARGIN_L, y=0.95, w=LEFT_W, row_h=0.28)

    # RIGHT SIDE — chart image
    chart_ph = slide_data.get("chart_placeholder", {})
    if chart_ph:
        _add_chart_image(slide, chart_ph.get("chart_key", "impact_ease_scatter"),
                         chart_paths, x=RIGHT_X, y=0.85, w=RIGHT_W)


def _render_low_hanging_fruit(prs, slide_data, layout_idx):
    """Slide 4: 3 easiest solutions — numbered blue titles + detail."""
    slide = prs.slides.add_slide(prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)])

    _title_block(slide, slide_data.get("title", "Low Hanging Fruit"))

    # Support both 'items' (new plan) and 'solutions' (legacy)
    items = slide_data.get("items", slide_data.get("solutions", []))

    box = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(0.95), Inches(CONTENT_W), Inches(6.0))
    tf = box.text_frame
    tf.word_wrap = True

    for idx, item in enumerate(items[:3], start=1):
        if isinstance(item, dict):
            action = _strip_md(str(item.get("action", item.get("title", ""))))
            detail = _strip_md(str(item.get("detail", "")))
            impact = _strip_md(str(item.get("impact", item.get("call_impact", ""))))
        else:
            action = _strip_md(str(item))
            detail = ""
            impact = ""

        # Numbered title (13pt, blue, bold)
        p = tf.paragraphs[0] if idx == 1 else tf.add_paragraph()
        r_title = p.add_run()
        r_title.text = f"{idx}. {action}"
        _apply(r_title, _style(13, bold=True, color=BLUE_ACCENT))

        # Detail line (11pt, black)
        if detail:
            p_det = tf.add_paragraph()
            r_det = p_det.add_run()
            r_det.text = f"    {detail}"
            _apply(r_det, _style(11, color=BLACK))

        # Impact line (10pt, muted)
        if impact:
            p_imp = tf.add_paragraph()
            r_imp = p_imp.add_run()
            r_imp.text = f"    {impact}"
            _apply(r_imp, _style(10, color=MID_GRAY))

        # Spacer
        if idx < min(len(items), 3):
            tf.add_paragraph()


def _render_recommendations(prs, slide_data, layout_idx):
    """Slide 5: 2x2 card grid with dimension accent bars."""
    slide = prs.slides.add_slide(prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)])

    _title_block(slide, slide_data.get("title", "Recommended Actions by Owning Team"))

    dimensions = slide_data.get("dimensions", [])
    if not dimensions:
        return

    # 2x2 grid filling 13.33" canvas
    CARD_W = 5.9
    CARD_H = 2.7
    CARD_GAP_X = CONTENT_W - 2 * CARD_W  # ~0.53
    CARD_GAP_Y = 0.3
    positions = [
        (MARGIN_L, 0.95),                                       # top-left
        (MARGIN_L + CARD_W + CARD_GAP_X, 0.95),                 # top-right
        (MARGIN_L, 0.95 + CARD_H + CARD_GAP_Y),                 # bottom-left
        (MARGIN_L + CARD_W + CARD_GAP_X, 0.95 + CARD_H + CARD_GAP_Y),  # bottom-right
    ]
    INNER_W = CARD_W - 0.4  # text area inside card

    for i, dim in enumerate(dimensions[:4]):
        if i >= len(positions):
            break
        cx, cy = positions[i]
        dim_name = str(dim.get("name", f"Dimension {i+1}"))
        accent = DIMENSION_COLORS.get(dim_name, BLUE_ACCENT)
        # Also check accent_color field if dimension name doesn't match
        if dim.get("accent_color"):
            accent = _hex_to_rgb(dim["accent_color"])

        # Card background
        _rect(slide, cx, cy, CARD_W, CARD_H, BG_LIGHT)
        # Left accent bar
        _rect(slide, cx, cy, 0.05, CARD_H, accent)

        # Dimension title: ■ Name
        box = slide.shapes.add_textbox(Inches(cx + 0.2), Inches(cy + 0.1), Inches(INNER_W), Inches(0.3))
        tf = box.text_frame
        tf.word_wrap = True
        r_sq = tf.paragraphs[0].add_run()
        r_sq.text = "\u25A0 "
        _apply(r_sq, _style(13, bold=True, color=accent))
        r_nm = tf.paragraphs[0].add_run()
        r_nm.text = dim_name
        _apply(r_nm, _style(13, bold=True, color=NAVY))

        # Action bullets
        actions = dim.get("actions", [])
        box = slide.shapes.add_textbox(Inches(cx + 0.2), Inches(cy + 0.5), Inches(INNER_W), Inches(CARD_H - 0.65))
        tf = box.text_frame
        tf.word_wrap = True
        for ai, act in enumerate(actions[:2]):
            act_title = _strip_md(str(act.get("title", "")))
            act_detail = _strip_md(str(act.get("detail", "")))

            p = tf.paragraphs[0] if ai == 0 else tf.add_paragraph()
            p.space_after = Pt(6)
            r_bullet = p.add_run()
            r_bullet.text = f"\u2022 {act_title}"
            _apply(r_bullet, _style(10.5, bold=True, color=BLACK))

            if act_detail:
                p2 = tf.add_paragraph()
                r_d = p2.add_run()
                r_d.text = f"   \u2192 {act_detail}"
                _apply(r_d, _style(9.5, color=MID_GRAY))


def _render_theme_card(prs, slide_data, layout_idx, chart_paths):
    """Slide 6+: Stats bar + LEFT narrative (60%) + RIGHT driver table (40%)."""
    slide = prs.slides.add_slide(prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)])

    title = _strip_md(slide_data.get("title", ""))
    stats_bar = slide_data.get("stats_bar", {})
    left_col = slide_data.get("left_column", {})
    right_col = slide_data.get("right_column", {})

    # Layout: LEFT 57% | gap | RIGHT 40%
    LEFT_W = 7.0
    RIGHT_X = MARGIN_L + LEFT_W + 0.5
    RIGHT_W = SLIDE_W - RIGHT_X - MARGIN_R  # ~4.83

    # 1. Theme title (22pt, NAVY)
    _tb(slide, MARGIN_L, 0.3, CONTENT_W, 0.35, title, SLIDE_TITLE)

    # 2. Stats bar (10pt, MID_GRAY)
    if stats_bar:
        parts = []
        if stats_bar.get("calls"):
            parts.append(f"Calls: {stats_bar['calls']}")
        if stats_bar.get("pct"):
            parts.append(stats_bar["pct"])
        if stats_bar.get("impact"):
            parts.append(f"Impact: {stats_bar['impact']}")
        if stats_bar.get("ease"):
            parts.append(f"Ease: {stats_bar['ease']}")
        if stats_bar.get("priority"):
            parts.append(f"Priority: {stats_bar['priority']}")
        _tb(slide, MARGIN_L, 0.62, CONTENT_W, 0.2, "  |  ".join(parts), STATS_TEXT)

    # 3. Rule
    _rule(slide, MARGIN_L, 0.82, CONTENT_W)

    # LEFT COLUMN
    if left_col:
        # CORE ISSUE (supports bullet points via newlines)
        _tb(slide, MARGIN_L, 0.95, LEFT_W, 0.25, "CORE ISSUE", H3_LABEL)
        core_issue = str(left_col.get("core_issue", ""))
        if core_issue:
            lines = [_strip_md(ln.strip()) for ln in core_issue.split("\n") if ln.strip()]
            if len(lines) > 1:
                box = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(1.22), Inches(LEFT_W), Inches(0.85))
                tf = box.text_frame
                tf.word_wrap = True
                for li, line in enumerate(lines[:4]):
                    p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                    p.space_after = Pt(2)
                    r = p.add_run()
                    r.text = line
                    _apply(r, BODY_TEXT)
            else:
                _tb(slide, MARGIN_L, 1.22, LEFT_W, 0.85, _strip_md(core_issue), BODY_TEXT)

        # PRIMARY DRIVER
        _tb(slide, MARGIN_L, 2.15, LEFT_W, 0.25, "PRIMARY DRIVER", H3_LABEL)
        driver = _strip_md(str(left_col.get("primary_driver", "")))
        if driver:
            _tb(slide, MARGIN_L, 2.42, LEFT_W, 0.55, driver, BODY_TEXT)

        # SOLUTIONS
        solutions = left_col.get("solutions", [])
        if solutions:
            _tb(slide, MARGIN_L, 3.05, LEFT_W, 0.25, "SOLUTIONS", H3_LABEL)

            box = slide.shapes.add_textbox(Inches(MARGIN_L), Inches(3.32), Inches(LEFT_W), Inches(3.7))
            tf = box.text_frame
            tf.word_wrap = True
            for si, sol in enumerate(solutions[:3]):
                action = _strip_md(str(sol.get("action", "")))
                dim = sol.get("dimension", "")
                classification = sol.get("classification", "")
                p = tf.paragraphs[0] if si == 0 else tf.add_paragraph()
                p.space_after = Pt(4)
                r_num = p.add_run()
                r_num.text = f"{si + 1}. "
                _apply(r_num, _style(11, bold=True, color=BLACK))
                r_act = p.add_run()
                r_act.text = f"{action} "
                _apply(r_act, _style(11, color=BLACK))
                if dim:
                    badge = f"[{dim}" + (f" · {classification}" if classification else "") + "]"
                    r_dim = p.add_run()
                    r_dim.text = badge
                    _apply(r_dim, _style(9, bold=True, color=BLUE_ACCENT))

    # RIGHT COLUMN — evidence/driver table
    if right_col:
        headers = right_col.get("headers", ["Driver", "Calls"])
        rows = right_col.get("rows", [])
        if headers and rows:
            n_cols = len(headers)
            if n_cols == 4:
                # Evidence table: Finding(55%) | Calls(10%) | Dimension(15%) | Action(20%)
                col_widths = [RIGHT_W * 0.55, RIGHT_W * 0.10, RIGHT_W * 0.15, RIGHT_W * 0.20]
            else:
                col_widths = [RIGHT_W - 1.3, 1.3]
            _add_table(slide, headers, rows, x=RIGHT_X, y=0.95, w=RIGHT_W,
                       col_widths=col_widths, row_h=0.28)


# ═══════════════════════════════════════════════════════════════════════════
# Legacy fallback renderer
# ═══════════════════════════════════════════════════════════════════════════

def _build_fallback_slide(prs, slide_data, layout_idx, chart_paths):
    """Generic fallback for slides with legacy `elements` arrays."""
    slide = prs.slides.add_slide(prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)])
    title = _strip_md(slide_data.get("title", ""))
    elements = slide_data.get("elements", [])

    if slide.shapes.title:
        slide.shapes.title.text = title

    body_ph = None
    for ph in slide.placeholders:
        if ph.placeholder_format.type in (2, 7):
            body_ph = ph
            break
    if body_ph and elements:
        tf = body_ph.text_frame
        tf.clear()
        first = True
        for elem in elements:
            etype = elem.get("type", "point_description")
            text = _strip_md(elem.get("text", ""))
            if not text or etype in ("table", "chart_placeholder"):
                continue
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            r = p.add_run()
            r.text = text
            _apply(r, BODY_TEXT)

    for elem in elements:
        if elem.get("type") == "chart_placeholder":
            _add_chart_image(slide, elem.get("chart_key", ""), chart_paths)


# ═══════════════════════════════════════════════════════════════════════════
# Public API
# ═══════════════════════════════════════════════════════════════════════════

def build_pptx_from_sections(
    section_blueprints: list[dict[str, Any]],
    chart_paths: dict[str, str],
    output_path: str,
    template_path: str = "",
    visual_hierarchy: dict[str, Any] | None = None,
) -> str:
    if template_path and Path(template_path).exists():
        prs = Presentation(template_path)
    else:
        prs = Presentation()

    # Canvas: 13.33 x 7.5 (widescreen 16:9)
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    n_layouts = len(prs.slide_layouts)

    for section in section_blueprints:
        for sd in section.get("slides", []):
            li = min(sd.get("layout_index", 1), n_layouts - 1)
            role = sd.get("slide_role", "content")

            if role in ("executive_summary", "hook", "hook_and_quick_wins"):
                _render_executive_summary(prs, sd, li)

            elif role == "pain_points" and sd.get("cards"):
                _render_pain_points(prs, sd, li)

            elif role in ("impact_ease", "impact_matrix"):
                _render_impact_ease(prs, sd, li, chart_paths)

            elif role == "low_hanging_fruit":
                _render_low_hanging_fruit(prs, sd, li)

            elif role == "biggest_bet":
                # Legacy: render as low hanging fruit fallback
                _render_low_hanging_fruit(prs, sd, li)

            elif role == "recommendations" and sd.get("dimensions"):
                _render_recommendations(prs, sd, li)

            elif role == "theme_card":
                _render_theme_card(prs, sd, li, chart_paths)

            else:
                _build_fallback_slide(prs, sd, li, chart_paths)

    prs.save(output_path)
    logger.info("PPTX built: %d slides -> %s",
                sum(len(s.get("slides", [])) for s in section_blueprints), output_path)
    return output_path
